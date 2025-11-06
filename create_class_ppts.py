#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
学生照片PPT生成工具
为每个班级创建PPT，每页显示一张学生照片及信息
"""

import os
import glob
import openpyxl
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

def load_students_by_class(excel_path):
    """从Excel文件中按班级加载学生信息"""
    print(f"正在读取Excel文件: {excel_path}")
    workbook = openpyxl.load_workbook(excel_path, read_only=True)
    students_by_class = {}  # {班级名: [(考号, 姓名), ...]}
    
    for sheet_name in workbook.sheetnames:
        print(f"正在处理sheet: {sheet_name}")
        students_by_class[sheet_name] = []
        sheet = workbook[sheet_name]
        
        for row in sheet.iter_rows(min_row=2, values_only=True):
            if len(row) >= 2 and row[0] and row[1]:
                exam_id, name = str(row[0]).strip(), str(row[1]).strip()
                if name and exam_id:
                    students_by_class[sheet_name].append((exam_id, name))
        
        # 按学号排序
        students_by_class[sheet_name].sort(key=lambda x: x[0])
        print(f"  {sheet_name}: {len(students_by_class[sheet_name])} 名学生")
    
    workbook.close()
    return students_by_class

def find_student_photos(directory):
    """查找所有学生照片"""
    pattern = os.path.join(directory, "*.png")
    photo_files = glob.glob(pattern)
    
    photos_dict = {}  # {(考号, 姓名): 照片路径}
    
    for photo_path in photo_files:
        filename = os.path.basename(photo_path)
        name_part = os.path.splitext(filename)[0]
        
        # 匹配模式：考号_姓名
        match = re.match(r'^(\d+)_(.+)$', name_part)
        if match:
            exam_id, name = match.groups()
            photos_dict[(exam_id, name)] = photo_path
    
    print(f"找到 {len(photos_dict)} 张学生照片")
    return photos_dict

def resize_image_for_ppt(image_path, max_width=8, max_height=6):
    """调整图片大小以适应PPT页面"""
    try:
        with Image.open(image_path) as img:
            # 获取原始尺寸
            orig_width, orig_height = img.size
            
            # 计算缩放比例
            width_ratio = (max_width * 96) / orig_width  # PPT中1英寸=96像素
            height_ratio = (max_height * 96) / orig_height
            scale_ratio = min(width_ratio, height_ratio, 1.0)  # 不放大，只缩小
            
            # 计算新尺寸（英寸）
            new_width = Inches(orig_width * scale_ratio / 96)
            new_height = Inches(orig_height * scale_ratio / 96)
            
            return new_width, new_height
    except Exception as e:
        print(f"处理图片 {image_path} 时出错: {e}")
        return Inches(6), Inches(4.5)  # 默认尺寸

def create_class_ppt(class_name, students, photos_dict, output_dir):
    """为指定班级创建PPT"""
    print(f"\n正在创建 {class_name} 的PPT...")
    
    # 创建新的PPT
    prs = Presentation()
    
    # 设置幻灯片尺寸（16:9）
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    students_with_photos = 0
    students_without_photos = 0
    
    for exam_id, name in students:
        # 添加新幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 检查是否有照片
        if (exam_id, name) in photos_dict:
            photo_path = photos_dict[(exam_id, name)]
            students_with_photos += 1
            
            try:
                # 计算图片尺寸
                img_width, img_height = resize_image_for_ppt(photo_path)
                
                # 添加图片（居中偏上）
                left = (prs.slide_width - img_width) / 2
                top = Inches(0.5)
                slide.shapes.add_picture(photo_path, left, top, img_width, img_height)
                
                # 添加学号和姓名文本框（图片下方）
                text_top = top + img_height + Inches(0.1)
                text_left = Inches(1)
                text_width = prs.slide_width - Inches(2)
                text_height = Inches(1.5)
                
                textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                text_frame = textbox.text_frame
                text_frame.clear()
                
                # 设置文本内容和格式
                p = text_frame.paragraphs[0]
                p.text = f"{exam_id}\n{name}"
                p.alignment = PP_ALIGN.CENTER
                
                # 设置字体格式：一号粗体红色
                font = p.font
                font.name = "微软雅黑"
                font.size = Pt(72)  # 一号字体大约72磅
                font.bold = True
                font.color.rgb = RGBColor(255, 0, 0)  # 红色
                
                print(f"  ✅ 已添加: {name} ({exam_id})")
                
            except Exception as e:
                print(f"  ❌ 添加 {name} ({exam_id}) 时出错: {e}")
                students_without_photos += 1
        else:
            students_without_photos += 1
            
            # 即使没有照片也创建一个页面，显示学生信息
            text_left = Inches(1)
            text_top = Inches(2.3)
            text_width = prs.slide_width - Inches(2)
            text_height = Inches(2)
            
            textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            p.text = f"{exam_id}\n{name}\n(无照片)"
            p.alignment = PP_ALIGN.CENTER
            
            font = p.font
            font.name = "微软雅黑"
            font.size = Pt(72)
            font.bold = True
            font.color.rgb = RGBColor(255, 0, 0)
            
            print(f"  ⚠️  无照片: {name} ({exam_id})")
    
    # 保存PPT
    ppt_filename = f"{class_name}_学生照片.pptx"
    ppt_path = os.path.join(output_dir, ppt_filename)
    prs.save(ppt_path)
    
    print(f"✅ {class_name} PPT 已保存: {ppt_path}")
    print(f"   共 {len(students)} 名学生，有照片 {students_with_photos} 人，无照片 {students_without_photos} 人")
    
    return ppt_path, students_with_photos, students_without_photos

def create_all_class_ppts(directory, excel_path):
    """为所有班级创建PPT的主函数"""
    print("="*60)
    print("学生照片PPT生成工具")
    print("="*60)
    
    # 加载学生信息
    students_by_class = load_students_by_class(excel_path)
    if not students_by_class:
        print("❌ 没有从Excel文件中读取到学生信息")
        return
    
    # 查找照片
    photos_dict = find_student_photos(directory)
    
    # 创建输出目录
    output_dir = os.path.join(directory, "班级PPT")
    os.makedirs(output_dir, exist_ok=True)
    print(f"\nPPT文件将保存到: {output_dir}")
    
    # 为每个班级创建PPT
    total_students = 0
    total_with_photos = 0
    total_without_photos = 0
    created_ppts = []
    
    for class_name, students in students_by_class.items():
        if students:  # 只处理有学生的班级
            ppt_path, with_photos, without_photos = create_class_ppt(
                class_name, students, photos_dict, output_dir
            )
            created_ppts.append(ppt_path)
            total_students += len(students)
            total_with_photos += with_photos
            total_without_photos += without_photos
    
    # 总结报告
    print("\n" + "="*60)
    print("PPT创建完成！")
    print("="*60)
    print(f"创建的PPT文件:")
    for ppt_path in created_ppts:
        print(f"  📄 {os.path.basename(ppt_path)}")
    
    print(f"\n统计信息:")
    print(f"  总学生数: {total_students}")
    print(f"  有照片: {total_with_photos}")
    print(f"  无照片: {total_without_photos}")
    print(f"  照片完成率: {(total_with_photos/total_students)*100:.1f}%")
    
    print(f"\n💡 PPT格式说明:")
    print(f"  - 每页显示一名学生")
    print(f"  - 照片居中显示")
    print(f"  - 学号和姓名显示在照片下方")
    print(f"  - 文字格式: 一号字体，粗体，红色")
    print(f"  - 页面按学号顺序排列")

def main():
    """主函数"""
    # 检查必要的库
    try:
        from pptx import Presentation
        from PIL import Image
    except ImportError as e:
        print(f"❌ 缺少必要的库: {e}")
        print("请安装必要的库:")
        print("pip install python-pptx pillow")
        return
    
    # 设置路径
    current_dir = os.getcwd()
    excel_path = os.path.join(current_dir, "mt2025.xlsx")
    
    # 检查Excel文件是否存在
    if not os.path.exists(excel_path):
        print(f"❌ Excel文件不存在: {excel_path}")
        return
    
    print(f"工作目录: {current_dir}")
    print(f"Excel文件: {excel_path}")
    
    # 创建PPT
    create_all_class_ppts(current_dir, excel_path)

if __name__ == "__main__":
    main()