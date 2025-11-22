#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
学生头像PPT生成工具（头像版）
为每个班级创建PPT，每页显示24张学生头像（6列×4行）
"""

import os
import glob
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def load_students_by_class(excel_path):
    """从Excel文件中按班级加载学生信息"""
    print(f"正在读取Excel文件: {excel_path}")
    workbook = openpyxl.load_workbook(excel_path, read_only=True)
    students_by_class = {}  # {班级名: [(考号, 姓名), ...]}
    
    for sheet_name in workbook.sheetnames:
        print(f"正在处理sheet: {sheet_name}")
        students_by_class[sheet_name] = []
        sheet = workbook[sheet_name]
        
        for row in sheet.iter_rows(min_row=2, values_only=True):
            if len(row) >= 2 and row[0] and row[1]:
                exam_id, name = str(row[0]).strip(), str(row[1]).strip()
                if name and exam_id:
                    students_by_class[sheet_name].append((exam_id, name))
        
        # 按学号排序
        students_by_class[sheet_name].sort(key=lambda x: x[0])
        print(f"  {sheet_name}: {len(students_by_class[sheet_name])} 名学生")
    
    workbook.close()
    return students_by_class


def find_headshot_photos(directory):
    """查找所有头像照片（文件名为9位考号）"""
    cuted_dir = os.path.join(directory, "cuted")
    if not os.path.exists(cuted_dir):
        print(f"❌ 头像目录不存在: {cuted_dir}")
        return {}
    
    pattern = os.path.join(cuted_dir, "*.png")
    photo_files = glob.glob(pattern)
    
    photos_dict = {}  # {考号: 照片路径}
    
    for photo_path in photo_files:
        filename = os.path.basename(photo_path)
        exam_id = os.path.splitext(filename)[0]
        
        # 验证是否为9位数字
        if exam_id.isdigit() and len(exam_id) == 9:
            photos_dict[exam_id] = photo_path
    
    print(f"在 {cuted_dir} 找到 {len(photos_dict)} 张头像照片")
    return photos_dict


def get_last_two_digits(exam_id):
    """获取考号的最后两位数字"""
    return exam_id[-2:] if len(exam_id) >= 2 else exam_id


def create_headshot_page(slide, students_page, photos_dict, prs, cols, rows):
    """在一页PPT上创建学生头像（根据cols和rows自动计算）"""
    # 页面设置由参数传入
    # cols: 列数
    # rows: 行数
    
    # 计算每个单元格的尺寸
    # 页面尺寸：13.33 x 7.5 英寸
    # 减小边距以获得更大的显示区域
    margin_left = Inches(0.2)
    margin_top = Inches(0.2)
    margin_right = Inches(0.2)
    margin_bottom = Inches(0.2)
    
    available_width = prs.slide_width - margin_left - margin_right
    available_height = prs.slide_height - margin_top - margin_bottom

    cell_width = available_width / cols
    cell_height = available_height / rows

    # 图片和文字的尺寸
    # 增大图片占比，让头像更清晰
    img_size = min(cell_width * 0.92, cell_height * 0.78)
    text_height = cell_height * 0.20  # 文字占20%高度
    
    students_with_photos = 0
    students_without_photos = 0
    
    for idx, (exam_id, name) in enumerate(students_page):
        if idx >= cols * rows:  # 最多 cols*rows 个
            break
        
        # 计算位置
        row = idx // cols
        col = idx % cols
        
        # 计算单元格中心位置
        cell_center_x = margin_left + col * cell_width + cell_width / 2
        cell_center_y = margin_top + row * cell_height + cell_height / 2
        
        # 图片位置（居中，偏上）
        img_left = cell_center_x - img_size / 2
        img_top = cell_center_y - img_size / 2 - text_height / 2
        
        # 添加头像
        if exam_id in photos_dict:
            photo_path = photos_dict[exam_id]
            students_with_photos += 1
            
            try:
                # 添加图片
                slide.shapes.add_picture(
                    photo_path, img_left, img_top, img_size, img_size
                )
            except Exception as e:
                print(f"  ❌ 添加 {name} ({exam_id}) 头像时出错: {e}")
                students_without_photos += 1
        else:
            students_without_photos += 1
            
            # 没有照片时，添加一个占位框
            shape = slide.shapes.add_shape(
                1,  # 矩形
                img_left, img_top, img_size, img_size
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # 浅灰色背景
            shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # 添加文字（考号后两位 + 姓名）
        text_left = cell_center_x - cell_width * 0.45
        text_top = img_top + img_size + Inches(0.05)
        text_width = cell_width * 0.9

        textbox = slide.shapes.add_textbox(
            text_left, text_top, text_width, text_height
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # 设置文本内容
        p = text_frame.paragraphs[0]
        last_two = get_last_two_digits(exam_id)
        p.text = f"{last_two}{name}"
        p.alignment = PP_ALIGN.CENTER
        
        # 设置字体格式：粗体红色
        font = p.font
        font.name = "微软雅黑"
        font.size = Pt(18)  # 适当的字体大小
        font.bold = True
        font.color.rgb = RGBColor(255, 0, 0)  # 红色

    return students_with_photos, students_without_photos


def create_class_headshot_ppt(
    class_name, students, photos_dict, output_dir, cols=6, rows=4
):
    """为指定班级创建头像版PPT"""
    print(f"\n正在创建 {class_name} 的头像PPT...")
    print(f"  布局: {cols}列 × {rows}行 (每页{cols * rows}人)")
    
    # 创建新的PPT
    prs = Presentation()
    
    # 根据行列比例自动设置幻灯片尺寸
    # 基础高度固定为7.5英寸，宽度根据列行比自动调整
    base_height = 7.5
    aspect_ratio = cols / rows  # 宽高比
    prs.slide_height = Inches(base_height)
    prs.slide_width = Inches(base_height * aspect_ratio)
    print(f"  页面尺寸: {base_height * aspect_ratio:.2f} × "
          f"{base_height} 英寸 (比例 {cols}:{rows})")
    
    total_with_photos = 0
    total_without_photos = 0
    
    # 每页学生数根据行列数自动计算
    students_per_page = cols * rows
    num_pages = (
        (len(students) + students_per_page - 1) // students_per_page
    )
    
    for page_idx in range(num_pages):
        # 添加新幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 获取本页的学生
        start_idx = page_idx * students_per_page
        end_idx = min(start_idx + students_per_page, len(students))
        students_page = students[start_idx:end_idx]

        # 创建本页内容
        with_photos, without_photos = create_headshot_page(
            slide, students_page, photos_dict, prs, cols, rows
        )
        total_with_photos += with_photos
        total_without_photos += without_photos
        
        print(f"  第 {page_idx + 1}/{num_pages} 页: {len(students_page)} 名学生")
    
    # 保存PPT
    ppt_filename = f"{class_name}_学生头像.pptx"
    ppt_path = os.path.join(output_dir, ppt_filename)
    prs.save(ppt_path)
    
    print(f"✅ {class_name} 头像PPT 已保存: {ppt_path}")
    print(
        f"   共 {len(students)} 名学生，"
        f"有头像 {total_with_photos} 人，"
        f"无头像 {total_without_photos} 人"
    )
    
    return ppt_path, total_with_photos, total_without_photos


def create_all_class_headshot_ppts(directory, excel_path, cols=6, rows=4):
    """为所有班级创建头像PPT的主函数"""
    print("="*60)
    print(f"学生头像PPT生成工具（布局：{cols}列×{rows}行）")
    print("="*60)
    
    # 加载学生信息
    students_by_class = load_students_by_class(excel_path)
    if not students_by_class:
        print("❌ 没有从Excel文件中读取到学生信息")
        return
    
    # 查找头像
    photos_dict = find_headshot_photos(directory)
    if not photos_dict:
        print("❌ 没有找到任何头像照片")
        return
    
    # 创建输出目录
    output_dir = os.path.join(directory, "班级PPT_头像版")
    os.makedirs(output_dir, exist_ok=True)
    print(f"\nPPT文件将保存到: {output_dir}")
    
    # 为每个班级创建PPT
    total_students = 0
    total_with_photos = 0
    total_without_photos = 0
    created_ppts = []
    
    for class_name, students in students_by_class.items():
        if students:  # 只处理有学生的班级
            ppt_path, with_photos, without_photos = (
                create_class_headshot_ppt(
                    class_name, students, photos_dict, output_dir, cols, rows
                )
            )
            created_ppts.append(ppt_path)
            total_students += len(students)
            total_with_photos += with_photos
            total_without_photos += without_photos
    
    # 总结报告
    print("\n" + "="*60)
    print("头像PPT创建完成！")
    print("="*60)
    print("创建的PPT文件:")
    for ppt_path in created_ppts:
        print(f"  📄 {os.path.basename(ppt_path)}")

    print("\n统计信息:")
    print(f"  总学生数: {total_students}")
    print(f"  有头像: {total_with_photos}")
    print(f"  无头像: {total_without_photos}")
    if total_students > 0:
        print(f"  头像完成率: {(total_with_photos/total_students)*100:.1f}%")

    print("\n💡 PPT格式说明:")
    print(f"  - 每页显示{cols * rows}名学生（{cols}列×{rows}行）")
    print("  - 使用cuted目录中的头像图片")
    print("  - 文字格式: 考号后两位+姓名，粗体，红色")
    print("  - 页面按学号顺序排列")


def main():
    """主函数"""
    # 检查必要的库
    try:
        from pptx import Presentation  # noqa: F401
        from PIL import Image  # noqa: F401
    except ImportError as e:
        print(f"❌ 缺少必要的库: {e}")
        print("请安装必要的库:")
        print("pip install python-pptx pillow")
        return
    
    # 设置路径
    current_dir = os.getcwd()
    excel_path = os.path.join(current_dir, "mt2025.xlsx")
    
    # 检查Excel文件是否存在
    if not os.path.exists(excel_path):
        print(f"❌ Excel文件不存在: {excel_path}")
        return
    
    # 检查cuted目录是否存在
    cuted_dir = os.path.join(current_dir, "cuted")
    if not os.path.exists(cuted_dir):
        print(f"❌ 头像目录不存在: {cuted_dir}")
        print("请创建cuted目录并放入学生头像照片（文件名为9位考号.png）")
        return
    
    print(f"工作目录: {current_dir}")
    print(f"Excel文件: {excel_path}")
    print(f"头像目录: {cuted_dir}")
    print("\n💡 如需修改布局，请编辑main()函数中的cols和rows参数")
    print("   默认: cols=6, rows=4 (每页24人)")
    print("   页面尺寸会根据行列比例自动调整\n")
    
    # 创建PPT - 可在此修改cols和rows参数调整布局
    # 例如: create_all_class_headshot_ppts(current_dir, excel_path, 5, 4)
    create_all_class_headshot_ppts(current_dir, excel_path, cols=5, rows=4)


if __name__ == "__main__":
    main()
